from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Pt
import argparse
import os
import json

initial_presentation = "./data/table_format.pptx"

def get_sundays_about_month(date_str):
    date = datetime.strptime(date_str, "%y.%m")
    first_day = date.replace(day=1)
    days_until_first_sunday = (6 - first_day.weekday()) % 7
    first_sunday = first_day + timedelta(days=days_until_first_sunday)
    
    sundays = []
    current_sunday = first_sunday
    while current_sunday.month == first_day.month:
        sundays.append(current_sunday.strftime("%y.%m.%d"))
        current_sunday += timedelta(days=7)
    
    return sundays

def get_formatted_date(start_date_str):
    start_date = datetime.strptime(start_date_str, "%y.%m.%d")
    end_date = start_date + timedelta(days=6)
    formatted_date_range = f"{start_date.month}월 {start_date.day}일 - {end_date.month}월 {end_date.day}일"
    
    return formatted_date_range

def get_weekly_plan(
        *, 
        start_book, 
        start_chapter, 
        daily_reading=2,
        bible_index="./data/bible_index.json", 
        bible_shortest_names="./data/bible_shortest_name.json"
    ):
    
    with open(bible_index, 'r') as file:
        bible_chapters = json.load(file)
    
    with open(bible_shortest_names, 'r') as file:
        shortest_names = json.load(file)
    
    start_chapter = int(start_chapter)
    weekly_plan = []
    current_book = start_book
    current_chapter = start_chapter

    for _ in range(7):
        chapters_left = bible_chapters[current_book] - current_chapter + 1
        if daily_reading <= chapters_left:
            end_chapter = current_chapter + daily_reading - 1
            reading_plan = f"{shortest_names[current_book]} {current_chapter} - {end_chapter}장"
            current_chapter += daily_reading
        else:
            next_book = next((book for book in bible_chapters if list(bible_chapters).index(book) > list(bible_chapters).index(current_book)), None)
            
            if next_book is None:
                next_book = list(bible_chapters.keys())[0]
                current_chapter = 1 + daily_reading - chapters_left
                reading_plan = f"{shortest_names[current_book]} {current_chapter} - {shortest_names[next_book]} {current_chapter - 1}장"
            else:
                end_chapter = daily_reading - chapters_left
                reading_plan = f"{shortest_names[current_book]} {current_chapter} - {shortest_names[next_book]} {end_chapter}장"
                current_book = next_book
                current_chapter = end_chapter + 1
            
        weekly_plan.append(reading_plan)
        
        if current_chapter > bible_chapters[current_book]:
            current_book = next((book for book in bible_chapters if list(bible_chapters).index(book) > list(bible_chapters).index(current_book)), None)
            if current_book is None:
                current_book = list(bible_chapters.keys())[0]
            current_chapter = 1
    
    return weekly_plan

def patch_cell(cell, text):
    cell.text = ""
    text_frame = cell.text_frame
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = text

    run.font.name = "배달의민족 한나체 Pro"
    run.font.size = Pt(46)

def patch_part_data(
        *,
        date,
        weekly_plan,
        selected_custom,
):
    pptx = Presentation(f"./custom/{selected_custom}.pptx")
    day_mapping = {
        "part_sun": 0,
        "part_mon": 1,
        "part_tue": 2,
        "part_wed": 3,
        "part_thu": 4,
        "part_fri": 5,
        "part_sat": 6,
    }

    for slide in pptx.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "part_date" in shape.text:
                    shape.text = shape.text.replace("part_date", f"( {date} )")

                for day_key, plan_index in day_mapping.items():
                    if day_key in shape.text:
                        shape.text = shape.text.replace(day_key, weekly_plan[plan_index])

    pptx.save(f"./shared/{date}.pptx")

    return
    


###############
#   main
###############

def create_weekly_plan(args):
    print(":: AUTO_WEEKLY_BIBLE_PLAN [ 25% ] - Calling functions...")
    start_book = args.start_book
    start_chapter = args.start_chapter
    date = get_formatted_date(args.date)
    selected_custom = args.selected_custom
    print(":: AUTO_WEEKLY_BIBLE_PLAN [ 50% ] - Preparing parameters...")

    weekly_plan = get_weekly_plan(start_book=start_book, start_chapter=start_chapter)
    print(":: AUTO_WEEKLY_BIBLE_PLAN [ 75% ] - Processing datas...")

    patch_part_data(date=date, weekly_plan=weekly_plan, selected_custom=selected_custom)
    print(":: AUTO_WEEKLY_BIBLE_PLAN [ 100% ] - Completed")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Create a weekly plan based on given parameters.")
    parser.add_argument("--start_book", required=True, help="The book to start with")
    parser.add_argument("--start_chapter", type=int, required=True, help="The chapter to start with")
    parser.add_argument("--date", required=True, help="The start date in YY.MM.DD format")
    parser.add_argument("--selected_custom", required=True, help="A custom selection for the plan")

    args = parser.parse_args()
    create_weekly_plan(args)